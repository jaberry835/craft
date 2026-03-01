"""
Grounding Service
Manages document grounding for agents using Azure AI Search and Azure Blob Storage.

This service enables RAG (Retrieval Augmented Generation) patterns where agents
can search through organizational documents stored in Azure Blob Storage.

Architecture (Sovereign Cloud Compatible):
- Documents are stored in Azure Blob Storage containers
- Documents are indexed into Azure AI Search with embeddings
- Agents get a custom search tool to query their grounded documents
- No dependency on Azure AI Foundry (works in Azure Government)
"""
from typing import Optional, Callable, Annotated
import asyncio
import hashlib
import io
import re
from datetime import datetime

from azure.storage.blob import ContainerClient, BlobClient
from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from azure.search.documents.indexes.models import (
    SearchIndex,
    SearchField,
    SearchFieldDataType,
    VectorSearch,
    HnswAlgorithmConfiguration,
    VectorSearchProfile,
    SearchableField,
    SimpleField
)
from azure.search.documents.models import VectorizedQuery

from config import get_settings, get_azure_credential
from observability import get_logger
from services.embedding_service import embedding_service
from services.security_token_service import security_token_service
from services.cosmos_service import cosmos_service
from services.document_intelligence_service import document_intelligence_service

settings = get_settings()
logger = get_logger(__name__)

# Index name prefix for grounding indices
GROUNDING_INDEX_PREFIX = "grounding-"


class GroundingService:
    """
    Manages document grounding for agents using Azure AI Search.
    
    When an agent is configured with grounding_sources (Azure Blob container URLs),
    this service:
    1. Creates a dedicated search index for that agent
    2. Indexes documents from the blob containers with embeddings
    3. Provides a search function that agents can use to query documents
    
    Works in sovereign clouds (Azure Government) - no Foundry dependency.
    """
    
    VECTOR_DIMENSIONS = 1536  # text-embedding-ada-002
    CHUNK_SIZE = 1000  # Characters per chunk
    CHUNK_OVERLAP = 200  # Overlap between chunks
    EXCEL_ROWS_PER_CHUNK = 20  # Rows per chunk for Excel (smaller = better entity-level precision)
    EMBEDDING_BATCH_SIZE = 16  # Chunks per embedding API call
    SEARCH_UPLOAD_BATCH_SIZE = 100  # Documents per search upload call
    MAX_CONCURRENT_BLOBS = 5  # Max blobs downloaded / parsed concurrently
    VECTOR_K_NEAREST = 50  # Vector candidates for RRF fusion (much larger than top to let keyword matches surface)

    # Regex for entity-like IDs (e.g. C001234, PO12345, INV-2024-001, TICK00042)
    # Matches 1-4 uppercase letters followed by optional separator then 3+ digits,
    # plus patterns with hyphens like INV-2024-001
    _ENTITY_ID_RE = re.compile(
        r'\b(?:[A-Z]{1,5}[-_]?\d{3,}(?:[-_]\d+)*)\b'
    )

    def __init__(self):
        self._credential = None
        self._index_client: Optional[SearchIndexClient] = None
        self._initialized = False
        self._lock = asyncio.Lock()
    
    async def initialize(self) -> None:
        """Initialize the grounding service with Azure AI Search."""
        if self._initialized:
            return
        
        async with self._lock:
            if self._initialized:
                return
            
            # Check if Azure AI Search is configured
            if not settings.search_endpoint:
                logger.warning(
                    "Azure AI Search endpoint not configured. "
                    "Grounding features will be disabled. "
                    "Set AZURE_SEARCH_ENDPOINT to enable document grounding."
                )
                self._initialized = True
                return
            
            try:
                # Use API key if available, otherwise managed identity
                if settings.search_key:
                    self._credential = AzureKeyCredential(settings.search_key)
                else:
                    self._credential = get_azure_credential()
                
                self._index_client = SearchIndexClient(
                    endpoint=settings.search_endpoint,
                    credential=self._credential
                )
                
                logger.info(f"Grounding service initialized with Azure AI Search: {settings.search_endpoint}")
                self._initialized = True
            except Exception as e:
                logger.error(f"Failed to initialize grounding service: {e}")
                self._initialized = True  # Mark as initialized to avoid retry loops
    
    @property
    def is_available(self) -> bool:
        """Check if grounding features are available."""
        return self._index_client is not None

    async def list_indexes(self) -> list[dict]:
        """List all Azure AI Search indexes on the configured service.
        
        Returns a list of dicts with index metadata (name, document count, etc.).
        Used by the admin UI to populate the "Use Existing Index" dropdown.
        """
        if not self.is_available:
            return []
        
        try:
            indexes = []
            for index in self._index_client.list_indexes():
                # Get document count via search client
                doc_count = None
                try:
                    sc = SearchClient(
                        endpoint=settings.search_endpoint,
                        index_name=index.name,
                        credential=self._credential
                    )
                    results = sc.search(
                        search_text="*",
                        select=["id"],
                        top=0,
                        include_total_count=True
                    )
                    doc_count = results.get_count()
                except Exception:
                    pass  # Count is best-effort
                
                indexes.append({
                    "name": index.name,
                    "field_count": len(index.fields) if index.fields else 0,
                    "document_count": doc_count,
                })
            
            logger.info(f"Listed {len(indexes)} indexes from Azure AI Search")
            return indexes
        except Exception as e:
            logger.error(f"Failed to list search indexes: {e}")
            return []

    def _get_index_name(self, agent_id: str) -> str:
        """Get the search index name for an agent's grounding documents."""
        # Create a safe index name from agent ID
        safe_id = agent_id.replace("-", "").lower()[:20]
        return f"{GROUNDING_INDEX_PREFIX}{safe_id}"
    
    async def _ensure_grounding_index(self, agent_id: str) -> str:
        """Ensure the grounding search index exists with the current schema.
        
        If the index exists but is missing required fields (e.g. ssToken),
        it is deleted and recreated so schema changes are picked up
        automatically without manual portal intervention.
        """
        index_name = self._get_index_name(agent_id)
        
        try:
            existing_index = self._index_client.get_index(index_name)
            # Check if schema is current – ssToken must be present
            existing_field_names = {f.name for f in existing_index.fields}
            if "ssToken" not in existing_field_names:
                logger.warning(f"Index {index_name} is missing 'ssToken' field – recreating index")
                self._index_client.delete_index(index_name)
                raise Exception("Schema outdated – recreate")
            logger.debug(f"Grounding index {index_name} exists with current schema")
        except Exception:
            # Create the index
            index = SearchIndex(
                name=index_name,
                fields=[
                    SimpleField(name="id", type=SearchFieldDataType.String, key=True),
                    SimpleField(name="agentId", type=SearchFieldDataType.String, filterable=True),
                    SimpleField(name="sourceUrl", type=SearchFieldDataType.String, filterable=True),
                    SimpleField(name="sourceName", type=SearchFieldDataType.String, filterable=True),
                    SimpleField(name="ssToken", type=SearchFieldDataType.String, filterable=True),
                    SearchableField(name="fileName", type=SearchFieldDataType.String),
                    SearchableField(name="content", type=SearchFieldDataType.String),
                    SimpleField(name="chunkIndex", type=SearchFieldDataType.Int32, sortable=True),
                    SimpleField(name="indexedAt", type=SearchFieldDataType.DateTimeOffset, sortable=True),
                    SearchField(
                        name="contentVector",
                        type=SearchFieldDataType.Collection(SearchFieldDataType.Single),
                        searchable=True,
                        vector_search_dimensions=self.VECTOR_DIMENSIONS,
                        vector_search_profile_name="vector-profile"
                    )
                ],
                vector_search=VectorSearch(
                    algorithms=[
                        HnswAlgorithmConfiguration(name="hnsw-config")
                    ],
                    profiles=[
                        VectorSearchProfile(
                            name="vector-profile",
                            algorithm_configuration_name="hnsw-config"
                        )
                    ]
                )
            )
            
            self._index_client.create_index(index)
            logger.info(f"Created grounding index: {index_name}")
        
        return index_name
    
    def _chunk_text(self, text: str) -> list[str]:
        """Split text into overlapping chunks."""
        chunks = []
        start = 0
        while start < len(text):
            end = start + self.CHUNK_SIZE
            chunk = text[start:end]
            if chunk.strip():
                chunks.append(chunk)
            start = end - self.CHUNK_OVERLAP
        return chunks

    def _chunk_tabular_text(self, text: str) -> list[str]:
        """
        Chunk tab-separated text (from Excel/CSV) preserving row boundaries
        and repeating column headers at the top of each chunk.

        Expected input format (from _extract_text_from_excel):
            --- Sheet: Customers ---
            CustomerID\tName\tSegment\tIndustry
            C001234\tContoso\tEnterprise\tTechnology
            ...

        Each chunk starts with the sheet header + column header row so the
        LLM can always interpret the data.
        """
        lines = text.split('\n')
        chunks: list[str] = []

        current_sheet_header = ''
        column_header_line = ''
        row_buffer: list[str] = []

        def _flush():
            """Write buffered rows as one chunk."""
            if not row_buffer:
                return
            header = ''
            if current_sheet_header:
                header = current_sheet_header + '\n'
            if column_header_line:
                header += column_header_line + '\n'
            chunks.append(header + '\n'.join(row_buffer))
            row_buffer.clear()

        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue

            # Detect sheet separators
            if stripped.startswith('--- Sheet:'):
                _flush()  # flush previous sheet
                current_sheet_header = stripped
                column_header_line = ''  # reset – next data row is header
                continue

            # First data row after a sheet header is treated as column header
            if current_sheet_header and not column_header_line:
                column_header_line = stripped
                continue

            row_buffer.append(stripped)

            if len(row_buffer) >= self.EXCEL_ROWS_PER_CHUNK:
                _flush()

        _flush()  # remaining rows
        return chunks
    
    async def _index_blob_documents(
        self,
        agent_id: str,
        index_name: str,
        container_url: str,
        source_name: str,
        blob_prefix: Optional[str] = None
    ) -> int:
        """
        Index documents from an Azure Blob Storage container.
        
        Uses batched embeddings and batched search uploads for performance.
        Returns the number of chunks indexed.
        """
        indexed_count = 0
        
        try:
            # Create container client with managed identity
            container_client = ContainerClient.from_container_url(
                container_url,
                credential=get_azure_credential()
            )
            
            # Create search client for this index
            search_client = SearchClient(
                endpoint=settings.search_endpoint,
                index_name=index_name,
                credential=self._credential
            )
            
            # List blobs (with optional prefix filter)
            blobs = list(container_client.list_blobs(name_starts_with=blob_prefix))
            
            # ---- Phase 1: Download & parse blobs concurrently ----
            semaphore = asyncio.Semaphore(self.MAX_CONCURRENT_BLOBS)
            
            async def _download_and_parse(blob):
                """Download a single blob and return (blob_name, ss_token, chunks) or None."""
                if not self._is_supported_file(blob.name):
                    logger.debug(f"Skipping unsupported file: {blob.name}")
                    return None
                try:
                    async with semaphore:
                        # Run sync blob I/O in a thread so we don't block the event loop
                        loop = asyncio.get_event_loop()
                        blob_client = container_client.get_blob_client(blob.name)
                        blob_properties = await loop.run_in_executor(None, blob_client.get_blob_properties)
                        metadata = blob_properties.metadata or {}
                        blob_ss_token = metadata.get("ss_tokens", "") or metadata.get("ss_token", "")

                        content = await loop.run_in_executor(
                            None, lambda: blob_client.download_blob().readall()
                        )
                    
                    # Parse content based on file type
                    # Try Document Intelligence first for supported binary files (especially PDFs)
                    text = None
                    if self._is_binary_document(blob.name) and document_intelligence_service.is_available and document_intelligence_service.supports_file(blob.name):
                        text = await document_intelligence_service.extract_text_or_none(content, blob.name)
                        if text:
                            logger.info(f"DI extracted {len(text)} chars from {blob.name}")
                    
                    # Fallback to local parsers for binary docs when DI is unavailable or failed
                    if text is None and self._is_binary_document(blob.name):
                        text = self._extract_text_from_binary(blob.name, content)
                        if not text:
                            logger.warning(f"No text extracted from binary document: {blob.name}")
                            return None
                    elif text is None:
                        try:
                            text = content.decode('utf-8')
                        except UnicodeDecodeError:
                            text = content.decode('latin-1')
                    
                    # Use tabular chunking for spreadsheet files, standard chunking otherwise
                    ext = self._get_ext(blob.name)
                    if ext in self._EXCEL_EXTENSIONS or ext in {'.csv', '.tsv'}:
                        chunks = self._chunk_tabular_text(text)
                    else:
                        chunks = self._chunk_text(text)
                    return (blob.name, blob_ss_token, chunks)
                except Exception as e:
                    logger.error(f"Failed to download/parse blob {blob.name}: {e}")
                    return None
            
            parsed_results = await asyncio.gather(
                *[_download_and_parse(b) for b in blobs]
            )
            
            # ---- Phase 2: Collect all chunks, then embed + upload incrementally ----
            # Build a flat list of (metadata, chunk_text) for batching
            chunk_records = []  # [(blob_name, ss_token, chunk_index, chunk_text), ...]
            for result in parsed_results:
                if result is None:
                    continue
                blob_name, blob_ss_token, chunks = result
                for i, chunk in enumerate(chunks):
                    chunk_records.append((blob_name, blob_ss_token, i, chunk))
            
            if not chunk_records:
                logger.info(f"No chunks to index from container {container_url}")
                return 0
            
            num_files = sum(1 for r in parsed_results if r)
            logger.info(f"Embedding {len(chunk_records)} chunks from {num_files} files")
            
            # Process in batches: embed → build docs → upload (incremental)
            upload_buffer: list[dict] = []
            for batch_start in range(0, len(chunk_records), self.EMBEDDING_BATCH_SIZE):
                batch_slice = chunk_records[batch_start:batch_start + self.EMBEDDING_BATCH_SIZE]
                batch_texts = [rec[3] for rec in batch_slice]
                batch_embeddings = await embedding_service.generate_embeddings(batch_texts)
                
                # Throttle to stay under TPM quota and avoid 429s
                await asyncio.sleep(0.2)
                
                # Build search documents for this batch
                for j, (blob_name, blob_ss_token, chunk_index, chunk_text) in enumerate(batch_slice):
                    if j >= len(batch_embeddings) or not batch_embeddings[j]:
                        logger.warning(f"Missing embedding for chunk {chunk_index} of {blob_name}")
                        continue
                    
                    doc_id = hashlib.md5(
                        f"{agent_id}:{container_url}:{blob_name}:{chunk_index}".encode()
                    ).hexdigest()
                    
                    upload_buffer.append({
                        "id": doc_id,
                        "agentId": agent_id,
                        "sourceUrl": container_url,
                        "sourceName": source_name,
                        "ssToken": blob_ss_token,
                        "fileName": blob_name,
                        "content": chunk_text,
                        "chunkIndex": chunk_index,
                        "indexedAt": datetime.utcnow().isoformat() + "Z",
                        "contentVector": batch_embeddings[j],
                    })
                
                # Flush upload buffer when full
                if len(upload_buffer) >= self.SEARCH_UPLOAD_BATCH_SIZE:
                    result = search_client.upload_documents(documents=upload_buffer)
                    succeeded = sum(1 for r in result if r.succeeded)
                    failed = len(upload_buffer) - succeeded
                    indexed_count += succeeded
                    if failed:
                        logger.error(f"Upload batch had {failed} failures out of {len(upload_buffer)} docs")
                        for r in result:
                            if not r.succeeded:
                                logger.error(f"  Doc {r.key}: {r.error_message}")
                    logger.info(f"Uploaded {indexed_count}/{len(chunk_records)} chunks to index")
                    upload_buffer = []
            
            # Flush remaining documents
            if upload_buffer:
                result = search_client.upload_documents(documents=upload_buffer)
                succeeded = sum(1 for r in result if r.succeeded)
                failed = len(upload_buffer) - succeeded
                indexed_count += succeeded
                if failed:
                    logger.error(f"Final upload batch had {failed} failures out of {len(upload_buffer)} docs")
                    for r in result:
                        if not r.succeeded:
                            logger.error(f"  Doc {r.key}: {r.error_message}")
            
            logger.info(f"Indexed {indexed_count} chunks from {num_files} files")
            
        except Exception as e:
            logger.error(f"Failed to access container {container_url}: {e}")
        
        return indexed_count
    
    # --- File type classification helpers ---

    _SUPPORTED_EXTENSIONS = {
        # Plain text / markup
        '.txt', '.md', '.mdx', '.rst', '.json', '.csv', '.tsv', '.xml',
        '.html', '.htm', '.log', '.yaml', '.yml', '.rtf',
        # Config files
        '.ini', '.cfg', '.conf', '.toml', '.env', '.properties',
        # Source code
        '.py', '.js', '.ts', '.java', '.cs', '.cpp', '.c', '.h',
        '.sql', '.sh', '.ps1', '.bat', '.cmd',
        # Binary document formats (need dedicated parsers)
        '.xlsx', '.xls', '.pdf', '.docx', '.pptx',
    }

    _EXCEL_EXTENSIONS = {'.xlsx', '.xls'}
    _PDF_EXTENSIONS = {'.pdf'}
    _DOCX_EXTENSIONS = {'.docx'}
    _PPTX_EXTENSIONS = {'.pptx'}

    def _get_ext(self, filename: str) -> str:
        """Return the lowercase extension of a filename (e.g. '.pdf')."""
        return '.' + filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

    def _is_supported_file(self, filename: str) -> bool:
        """Check if the file type is supported for indexing."""
        return self._get_ext(filename) in self._SUPPORTED_EXTENSIONS

    def _is_binary_document(self, filename: str) -> bool:
        """Return True for file types that need a specialised parser."""
        ext = self._get_ext(filename)
        return ext in (self._EXCEL_EXTENSIONS | self._PDF_EXTENSIONS
                       | self._DOCX_EXTENSIONS | self._PPTX_EXTENSIONS)

    # --- Binary document text extraction ---

    def _extract_text_from_binary(self, filename: str, content: bytes) -> str:
        """Dispatch to the correct parser for binary document types."""
        ext = self._get_ext(filename)
        if ext in self._EXCEL_EXTENSIONS:
            return self._extract_text_from_excel(content)
        if ext in self._PDF_EXTENSIONS:
            return self._extract_text_from_pdf(content)
        if ext in self._DOCX_EXTENSIONS:
            return self._extract_text_from_docx(content)
        if ext in self._PPTX_EXTENSIONS:
            return self._extract_text_from_pptx(content)
        return ''

    def _extract_text_from_excel(self, content: bytes) -> str:
        """
        Extract text content from an Excel file (.xlsx/.xls) using openpyxl.
        
        - Filters out cells with formula errors (#REF!, #N/A, #VALUE!, etc.)
        - Returns structured text with clear sheet/header context so that
          downstream chunking can preserve column headers per chunk.
        """
        # Patterns that indicate broken formula references
        _ERROR_TOKENS = {'#REF!', '#N/A', '#VALUE!', '#DIV/0!', '#NAME?', '#NULL!', '#NUM!'}

        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    row_values = []
                    for cell in row:
                        if cell is None:
                            row_values.append('')
                        else:
                            s = str(cell).strip()
                            # Replace Excel error tokens with empty string
                            if s in _ERROR_TOKENS:
                                row_values.append('')
                            else:
                                row_values.append(s)
                    # Skip completely empty rows
                    if any(v for v in row_values):
                        text_parts.append('\t'.join(row_values))
            wb.close()
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract text from Excel file: {e}")
            return ''

    def _extract_text_from_pdf(self, content: bytes) -> str:
        """Extract text from a PDF using pymupdf (fitz)."""
        try:
            import pymupdf
            doc = pymupdf.open(stream=content, filetype="pdf")
            text_parts = []
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"--- Page {page_num} ---")
                    text_parts.append(text)
            doc.close()
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract text from PDF: {e}")
            return ''

    def _extract_text_from_docx(self, content: bytes) -> str:
        """Extract text from a Word .docx file using python-docx."""
        try:
            import docx
            doc = docx.Document(io.BytesIO(content))
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_values = [cell.text.strip() for cell in row.cells]
                    if any(row_values):
                        text_parts.append('\t'.join(row_values))
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract text from DOCX: {e}")
            return ''

    def _extract_text_from_pptx(self, content: bytes) -> str:
        """Extract text from a PowerPoint .pptx file using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_texts.append(para.text)
                if slide_texts:
                    text_parts.append(f"--- Slide {slide_num} ---")
                    text_parts.extend(slide_texts)
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract text from PPTX: {e}")
            return ''
    
    async def create_or_update_grounding_index(
        self,
        agent_id: str,
        agent_name: str,
        grounding_sources: list[dict]
    ) -> Optional[str]:
        """
        Create or update the grounding index for an agent.
        
        Args:
            agent_id: Unique identifier for the agent
            agent_name: Human-readable agent name
            grounding_sources: List of grounding source configurations
            
        Returns:
            Index name if successful, None otherwise
        """
        if not self.is_available:
            logger.warning("Grounding service not available - skipping index creation")
            return None
        
        if not grounding_sources:
            return None
        
        try:
            # Ensure index exists
            index_name = await self._ensure_grounding_index(agent_id)
            
            # Clear existing documents for this agent (to handle updates)
            await self._clear_agent_documents(index_name, agent_id)
            
            # Index documents from each grounding source
            total_indexed = 0
            for source in grounding_sources:
                container_url = source.get("container_url", "")
                source_name = source.get("name") or self._extract_container_name(container_url)
                blob_prefix = source.get("blob_prefix")
                
                count = await self._index_blob_documents(
                    agent_id=agent_id,
                    index_name=index_name,
                    container_url=container_url,
                    source_name=source_name,
                    blob_prefix=blob_prefix
                )
                total_indexed += count
            
            logger.info(f"Indexed {total_indexed} total chunks for agent {agent_id} in index {index_name}")
            return index_name
            
        except Exception as e:
            logger.error(f"Failed to create grounding index for agent {agent_id}: {e}")
            return None
    
    async def _clear_agent_documents(self, index_name: str, agent_id: str) -> None:
        """Clear all documents for an agent from the index."""
        try:
            search_client = SearchClient(
                endpoint=settings.search_endpoint,
                index_name=index_name,
                credential=self._credential
            )
            
            # Find all documents for this agent
            results = search_client.search(
                search_text="*",
                filter=f"agentId eq '{agent_id}'",
                select=["id"],
                top=1000
            )
            
            doc_ids = [{"id": r["id"]} for r in results]
            if doc_ids:
                search_client.delete_documents(documents=doc_ids)
                logger.info(f"Cleared {len(doc_ids)} existing documents for agent {agent_id}")
        except Exception as e:
            logger.debug(f"Could not clear documents (index may not exist yet): {e}")
    
    async def delete_grounding_index(self, agent_id: str) -> bool:
        """Delete the grounding index for an agent."""
        if not self.is_available:
            return False
        
        try:
            index_name = self._get_index_name(agent_id)
            self._index_client.delete_index(index_name)
            logger.info(f"Deleted grounding index: {index_name}")
            return True
        except Exception as e:
            logger.warning(f"Failed to delete grounding index for {agent_id}: {e}")
            return False
    
    def _extract_container_name(self, url: str) -> str:
        """Extract container name from URL."""
        try:
            parts = url.rstrip('/').split('/')
            return parts[-1] if parts else 'documents'
        except:
            return 'documents'
    
    async def search_grounding_documents(
        self,
        agent_id: str,
        query: str,
        top_k: int = 5,
        user_token: Optional[str] = None,
        index_name_override: Optional[str] = None
    ) -> list[dict]:
        """
        Search grounding documents for an agent with security filtering.
        
        Args:
            agent_id: The agent ID
            query: Search query text
            top_k: Number of results to return
            user_token: User's bearer token for SS token resolution.
                        If None or access checker is not configured, no
                        security filtering is applied (all docs returned).
            index_name_override: If set, use this index name instead of
                                 deriving it from agent_id. Used for
                                 external (BYOI) indexes.
            
        Returns:
            List of matching document chunks with content and metadata
        """
        if not self.is_available:
            return []
        
        try:
            index_name = index_name_override or self._get_index_name(agent_id)
            is_external = bool(index_name_override)
            if is_external:
                logger.info(f"Using external (BYOI) index: {index_name} for agent {agent_id}")
            
            # Generate query embedding
            query_embedding = await embedding_service.generate_embedding(query)
            if not query_embedding:
                logger.warning("Failed to generate query embedding")
                return []
            
            search_client = SearchClient(
                endpoint=settings.search_endpoint,
                index_name=index_name,
                credential=self._credential
            )
            
            vector_query = VectorizedQuery(
                vector=query_embedding,
                k_nearest_neighbors=self.VECTOR_K_NEAREST,
                fields="contentVector"
            )
            
            # Build filter: for managed indexes filter by agentId;
            # for external (BYOI) indexes skip the agentId filter since the
            # documents were indexed externally and may have a different (or no)
            # agentId value.
            filter_parts = []
            if not is_external:
                filter_parts.append(f"agentId eq '{agent_id}'")
            
            if user_token and security_token_service.is_available:
                ss_tokens = await security_token_service.get_user_ss_tokens(user_token)
                if ss_tokens is not None:
                    security_filter = security_token_service.build_security_filter(ss_tokens)
                    if security_filter:
                        filter_parts.append(security_filter)
                        logger.info(f"Security filter applied: {len(ss_tokens)} SS tokens")
                    else:
                        # User has zero tokens – they have no access to any documents
                        logger.warning("User has no SS tokens – returning empty results")
                        return []
                # If ss_tokens is None the API call failed – fall through without filter
                # (fail-open; change to fail-closed by returning [] if desired)
            
            odata_filter = " and ".join(filter_parts) if filter_parts else None
            
            select_fields = ["id", "fileName", "content", "sourceName", "chunkIndex"]

            # -- Pass 1: Hybrid search (vector + BM25 via RRF fusion) --
            hybrid_results = search_client.search(
                search_text=query,
                vector_queries=[vector_query],
                filter=odata_filter,
                select=select_fields,
                top=top_k
            )

            # -- Pass 2: Pure keyword / BM25 search (no vector) --
            # This guarantees exact keyword matches (e.g. opaque IDs like
            # C001234) always surface, even when they have no semantic
            # similarity to the query embedding.
            keyword_results = search_client.search(
                search_text=query,
                filter=odata_filter,
                select=select_fields,
                top=top_k
            )

            # -- Pass 3: Entity ID exact search --
            # If the query contains entity-like IDs (C001234, PO12345, etc.),
            # do a focused search with JUST the ID.  This prevents generic
            # words like "customer", "segment", "region" from drowning out
            # the exact ID match in BM25 scoring.
            entity_ids = self._ENTITY_ID_RE.findall(query.upper())
            # Also check original case for mixed-case queries
            entity_ids += self._ENTITY_ID_RE.findall(query)
            entity_ids = list(dict.fromkeys(entity_ids))  # dedupe, preserve order

            entity_results_list = []
            for eid in entity_ids[:3]:  # max 3 entity ID searches
                logger.info(f"Entity ID search pass for: '{eid}'")
                eid_results = search_client.search(
                    search_text=eid,
                    filter=odata_filter,
                    select=select_fields,
                    top=top_k
                )
                entity_results_list.append(eid_results)

            # -- Merge & deduplicate, keeping the higher score per doc --
            seen: dict[str, dict] = {}  # id -> doc dict
            for result in hybrid_results:
                doc_id = result["id"]
                seen[doc_id] = {
                    "file_name": result.get("fileName", ""),
                    "source": result.get("sourceName", ""),
                    "content": result.get("content", ""),
                    "chunk_index": result.get("chunkIndex", 0),
                    "score": result.get("@search.score", 0)
                }
            for result in keyword_results:
                doc_id = result["id"]
                score = result.get("@search.score", 0)
                if doc_id not in seen or score > seen[doc_id]["score"]:
                    seen[doc_id] = {
                        "file_name": result.get("fileName", ""),
                        "source": result.get("sourceName", ""),
                        "content": result.get("content", ""),
                        "chunk_index": result.get("chunkIndex", 0),
                        "score": score
                    }
            # Merge entity ID results — boost score ONLY when the chunk
            # actually contains the entity ID.  This avoids inflating
            # tangentially-related chunks that BM25 returned on partial
            # token overlap.
            ENTITY_BOOST = 100.0
            for eid, eid_results in zip(entity_ids[:3], entity_results_list):
                eid_upper = eid.upper()
                for result in eid_results:
                    doc_id = result["id"]
                    raw_score = result.get("@search.score", 0)
                    content = result.get("content", "")
                    # Only boost if the entity ID literally appears in the chunk
                    has_entity = eid_upper in content.upper()
                    boosted_score = raw_score + ENTITY_BOOST if has_entity else raw_score
                    if doc_id not in seen or boosted_score > seen[doc_id]["score"]:
                        seen[doc_id] = {
                            "file_name": result.get("fileName", ""),
                            "source": result.get("sourceName", ""),
                            "content": content,
                            "chunk_index": result.get("chunkIndex", 0),
                            "score": boosted_score
                        }

            # Sort by score descending, return top_k
            documents = sorted(seen.values(), key=lambda d: d["score"], reverse=True)[:top_k]

            # Log results with content snippets at INFO for visibility
            if documents:
                logger.info(f"Grounding search for agent {agent_id}: query='{query[:80]}', "
                            f"results={len(documents)}, top_score={documents[0]['score']:.4f}, "
                            f"hybrid_count={len([d for d in seen.values()])}, "
                            f"unique_after_merge={len(seen)}")
                for i, d in enumerate(documents[:5]):
                    logger.info(f"  result#{i+1} score={d['score']:.4f} file={d['file_name']} "
                                f"chunk={d['chunk_index']}: {d['content'][:150]}...")
            else:
                logger.info(f"Grounding search for agent {agent_id}: query='{query[:80]}', results=0")

            return documents
            
        except Exception as e:
            logger.error(f"Failed to search grounding documents for agent {agent_id}: {e}")
            return []
    
    def create_search_tool(
        self, agent_id: str, agent_name: str,
        user_token: Optional[str] = None,
        index_name_override: Optional[str] = None
    ) -> Callable:
        """
        Create a search tool function that an agent can use to query its grounded documents.
        
        This returns a function that can be passed to the agent as a tool.
        The user_token is captured in the closure so security filtering is applied
        transparently on every search call.
        
        Args:
            agent_id: The agent's unique ID
            agent_name: The agent's display name
            user_token: Bearer token for SS token security filtering
            index_name_override: For external (BYOI) indexes, the index name to query
        """
        # Capture user_token in closure for security filtering
        _user_token = user_token
        _index_override = index_name_override
        
        async def search_knowledge_base(
            query: Annotated[str, "The search query to find relevant information in the agent's knowledge base documents"]
        ) -> str:
            """
            Search the agent's knowledge base for relevant information.
            Use this tool when you need to find specific information from your grounding documents.
            Always reference the source file names in your response so users know
            which documents the information came from.
            """
            results = await self.search_grounding_documents(
                agent_id, query, top_k=10,
                user_token=_user_token,
                index_name_override=_index_override
            )
            
            if not results:
                return "No relevant documents found in the knowledge base."
            
            # Format results for the agent with source attributions
            # NOTE: No URLs here — the frontend auto-links file names to the
            # document viewer endpoint, which is more reliable than asking the
            # LLM to preserve URLs.
            formatted = []
            seen_files = set()
            
            for doc in results:
                file_name = doc['file_name']
                source = doc['source']
                seen_files.add(file_name)
                
                formatted.append(
                    f"**Source: {source} — {file_name}**\n"
                    f"{doc['content']}\n"
                )
            
            result_text = "\n---\n".join(formatted)
            
            # Append a simple sources list for easy reference
            if seen_files:
                sources_list = ", ".join(sorted(seen_files))
                result_text += f"\n\nDocuments referenced: {sources_list}"
            
            return result_text
        
        return search_knowledge_base
    
    async def validate_container_access(self, container_url: str) -> tuple[bool, str]:
        """
        Validate that the container URL is accessible.
        
        Args:
            container_url: Azure Blob Storage container URL
            
        Returns:
            Tuple of (success: bool, message: str)
        """
        # Basic URL validation
        if not container_url.startswith("https://") or ".blob." not in container_url:
            return False, "Invalid Azure Blob Storage URL format. Expected: https://<account>.blob.core.windows.net/<container>"
        
        try:
            # Try to list blobs (just check access)
            container_client = ContainerClient.from_container_url(
                container_url,
                credential=get_azure_credential()
            )
            # Just try to get container properties to verify access
            container_client.get_container_properties()
            return True, "Container is accessible"
        except Exception as e:
            return False, f"Cannot access container: {str(e)}"
    
    async def get_index_status(self, agent_id: str) -> Optional[dict]:
        """
        Get the status of a grounding index.
        
        Args:
            agent_id: The agent ID
            
        Returns:
            Status dict with document count, etc. or None if not found
        """
        if not self.is_available:
            return None
        
        try:
            index_name = self._get_index_name(agent_id)
            
            search_client = SearchClient(
                endpoint=settings.search_endpoint,
                index_name=index_name,
                credential=self._credential
            )
            
            # Count documents for this agent
            results = search_client.search(
                search_text="*",
                filter=f"agentId eq '{agent_id}'",
                select=["id"],
                top=0,
                include_total_count=True
            )
            
            return {
                "index_name": index_name,
                "document_count": results.get_count() or 0,
                "status": "ready"
            }
        except Exception as e:
            logger.debug(f"Could not get index status for {agent_id}: {e}")
            return None


    async def fetch_blob_content(self, agent_id: str, file_name: str) -> Optional[tuple[str, str, str]]:
        """
        Fetch a blob's raw content from the agent's grounding sources.
        
        Looks up the agent's configured grounding sources, finds the blob,
        and returns its content along with its SS token for authorization checks.
        
        Args:
            agent_id: The agent ID
            file_name: The blob file name
            
        Returns:
            Tuple of (content_text, content_type, ss_token) or None if not found.
            ss_token will be empty string if the blob has no security token.
        """
        try:
            agent = await cosmos_service.get_agent(agent_id)
            if not agent:
                return None
            
            grounding_sources = agent.get("grounding_sources", [])
            
            for source in grounding_sources:
                container_url = source.get("container_url", "")
                if not container_url:
                    continue
                
                try:
                    container_client = ContainerClient.from_container_url(
                        container_url,
                        credential=get_azure_credential()
                    )
                    blob_client = container_client.get_blob_client(file_name)
                    
                    # Get blob properties to read metadata (ss_token)
                    blob_properties = blob_client.get_blob_properties()
                    metadata = blob_properties.metadata or {}
                    blob_ss_token = metadata.get("ss_tokens", "") or metadata.get("ss_token", "")
                    
                    download = blob_client.download_blob()
                    content = download.readall()
                    
                    # Determine content type from extension
                    ext = ('.' + file_name.rsplit('.', 1)[-1].lower()) if '.' in file_name else ''
                    content_type_map = {
                        '.md': 'text/markdown; charset=utf-8',
                        '.txt': 'text/plain; charset=utf-8',
                        '.json': 'application/json; charset=utf-8',
                        '.csv': 'text/csv; charset=utf-8',
                        '.pdf': 'application/pdf',
                    }
                    content_type = content_type_map.get(ext, 'text/plain; charset=utf-8')
                    
                    # Decode text content
                    try:
                        text = content.decode('utf-8')
                    except UnicodeDecodeError:
                        text = content.decode('latin-1')
                    
                    return text, content_type, blob_ss_token
                except Exception:
                    # Blob not in this source, try next
                    continue
            
            return None
        except Exception as e:
            logger.error(f"Failed to fetch blob {file_name} for agent {agent_id}: {e}")
            return None

    async def fetch_blob_by_url(self, blob_url: str) -> Optional[tuple[bytes, str, str]]:
        """
        Fetch a blob directly by its full URL using managed identity.
        
        Used to proxy blobs referenced by MCP tools or other external sources
        without requiring SAS tokens.
        
        Args:
            blob_url: Full Azure Blob Storage URL
                      (e.g. https://account.blob.core.usgovcloudapi.net/container/file.txt)
            
        Returns:
            Tuple of (raw_bytes, content_type, ss_token) or None if not found.
            ss_token is empty string if the blob has no security token.
        """
        try:
            blob_client = BlobClient.from_blob_url(
                blob_url,
                credential=get_azure_credential()
            )
            
            # Get blob properties for metadata (ss_token)
            blob_properties = blob_client.get_blob_properties()
            metadata = blob_properties.metadata or {}
            blob_ss_token = metadata.get("ss_tokens", "") or metadata.get("ss_token", "")
            
            # Determine content type from extension
            blob_name = blob_client.blob_name
            ext = ('.' + blob_name.rsplit('.', 1)[-1].lower()) if '.' in blob_name else ''
            content_type_map = {
                '.md': 'text/markdown; charset=utf-8',
                '.txt': 'text/plain; charset=utf-8',
                '.json': 'application/json; charset=utf-8',
                '.csv': 'text/csv; charset=utf-8',
                '.pdf': 'application/pdf',
                '.html': 'text/html; charset=utf-8',
                '.xml': 'text/xml; charset=utf-8',
            }
            content_type = content_type_map.get(ext, 'application/octet-stream')
            
            download = blob_client.download_blob()
            content = download.readall()
            
            return content, content_type, blob_ss_token
        except Exception as e:
            logger.error(f"Failed to fetch blob by URL {blob_url}: {e}")
            return None


# Global singleton instance
grounding_service = GroundingService()
